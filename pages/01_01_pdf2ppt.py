import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image

st.title("PDF를 이미지로 변환하여 PPTX로 저장하기 📄➡️📽️")
st.write("PDF 파일을 업로드하면 각 페이지를 이미지로 변환하고, 이를 PPTX 파일로 만들어 다운로드할 수 있어요! 😊")

uploaded_pdf = st.file_uploader("PDF 파일을 업로드 해주세요.", type=["pdf"])

if uploaded_pdf is not None:
    pdf_bytes = uploaded_pdf.read()
    st.info("PDF 파일을 이미지로 변환하는 중입니다... 🕒")

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    st.success(f"PDF에서 {len(doc)} 페이지를 이미지로 변환했어요! 🎉")

    # 프레젠테이션 생성
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)  # 해상도 조절 가능
        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        slide = prs.slides.add_slide(blank_slide_layout)
        image_stream = BytesIO()
        image.save(image_stream, format="PNG")
        image_stream.seek(0)

        width = pix.width / 150  # dpi로 나눠서 인치로 계산
        height = pix.height / 150

        slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=Inches(width), height=Inches(height))
        st.write(f"페이지 {page_num + 1} 변환 완료! 🖼️")

    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)

    st.success("모든 페이지를 PPTX 파일로 변환했어요! 🎯")
    st.download_button(
        label="PPTX 파일 다운로드 ⬇️",
        data=pptx_bytes,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
else:
    st.warning("PDF 파일을 업로드해주세요! ☝️")
